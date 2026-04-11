#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "presentation_assets"
OUT = ROOT / "TOPIC1_PRESENTATION_DRAFT_V2.pptx"

TITLE_COLOR = RGBColor(32, 45, 64)
TEXT_COLOR = RGBColor(45, 45, 45)
ACCENT_BLUE = RGBColor(76, 120, 168)
ACCENT_ORANGE = RGBColor(245, 133, 24)
ACCENT_GREEN = RGBColor(84, 162, 75)
LIGHT_BG = RGBColor(246, 248, 251)
LIGHT_GRAY = RGBColor(110, 118, 129)


def add_textbox(slide, left, top, width, height, text="", font_size=20, bold=False,
                color=TEXT_COLOR, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_bullets(slide, left, top, width, height, bullets, font_size=20, color=TEXT_COLOR,
                level0_space=10, level1_font_size=None):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            level, text = item
        else:
            level, text = 0, item
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.space_after = Pt(level0_space if level == 0 else 4)
        run = p.runs[0]
        run.font.size = Pt(level1_font_size if level == 1 and level1_font_size else font_size)
        run.font.color.rgb = color
    return box


def add_title(slide, title):
    add_textbox(slide, Inches(0.55), Inches(0.22), Inches(11.8), Inches(0.5),
                title, font_size=24, bold=True, color=TITLE_COLOR)


def add_footer(slide, text="Draft generated automatically from project results"):
    add_textbox(slide, Inches(0.55), Inches(6.95), Inches(12.0), Inches(0.2),
                text, font_size=9, color=LIGHT_GRAY)


def add_header_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def add_note_box(slide, left, top, width, height, title, body_lines):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = RGBColor(220, 225, 232)
    add_textbox(slide, left + Inches(0.12), top + Inches(0.05), width - Inches(0.24), Inches(0.2),
                title, font_size=12, bold=True, color=TITLE_COLOR)
    add_bullets(
        slide,
        left + Inches(0.12),
        top + Inches(0.3),
        width - Inches(0.24),
        height - Inches(0.35),
        body_lines,
        font_size=11,
        level1_font_size=10,
    )


def slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title(slide, "Main result snapshot: post-training changed different metrics, not all benchmarks uniformly")
    slide.shapes.add_picture(str(ASSETS / "chart_slide1_benchmark_snapshot_v2.png"), Inches(6.55), Inches(0.95), width=Inches(6.15))
    bullets = [
        "On GSM8K XML, GRPO leads on accuracy and strict_accuracy.",
        "On GSM8K XML, SFT leads on content_accuracy.",
        "On ARC-Easy, there is no general-capability gain after post-training.",
        "So post-training changed targeted behaviors, not every metric in the same direction.",
        "The gains are real, but they are benchmark-specific and metric-specific.",
    ]
    add_bullets(slide, Inches(0.7), Inches(1.2), Inches(5.4), Inches(3.8), bullets, font_size=20)
    add_note_box(
        slide,
        Inches(0.72),
        Inches(5.3),
        Inches(5.1),
        Inches(1.15),
        "Main takeaway",
        [
            "GRPO's strongest gain is answer landing under the XML protocol.",
            "SFT's strongest gain is getting correct answer content into the output.",
            "ARC-Easy does not show a general-capability improvement.",
        ],
    )
    add_footer(slide)


def slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title(slide, "Main benchmark: GSM8K XML lets us separate content generation from answer landing")

    setup = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.95), Inches(11.95), Inches(0.45))
    setup.fill.solid()
    setup.fill.fore_color.rgb = LIGHT_BG
    setup.line.color.rgb = RGBColor(220, 225, 232)
    add_textbox(
        slide, Inches(0.9), Inches(1.03), Inches(11.5), Inches(0.2),
        "GSM8K test set | limit=300 | max_new_tokens=1024 | XML prompt | greedy decoding",
        font_size=14, bold=True, color=TITLE_COLOR, align=PP_ALIGN.CENTER
    )

    slide.shapes.add_picture(str(ASSETS / "chart_gsm8k_core_metrics.png"), Inches(0.7), Inches(1.55), width=Inches(7.4))

    add_note_box(
        slide,
        Inches(8.35),
        Inches(1.55),
        Inches(4.2),
        Inches(2.35),
        "How to read the metrics",
        [
            "accuracy: final answer judged correct",
            "strict_accuracy: answer is explicitly and cleanly extractable",
            "content_accuracy: correct answer appears in the output content",
        ],
    )
    add_note_box(
        slide,
        Inches(8.35),
        Inches(4.05),
        Inches(4.2),
        Inches(1.55),
        "Reading note",
        [
            "GSM8K XML is a diagnostic benchmark, not a pure math score.",
            "It mixes content correctness, answer landing, format extractability, and truncation.",
        ],
    )
    add_textbox(
        slide,
        Inches(0.75),
        Inches(5.95),
        Inches(11.9),
        Inches(0.4),
        "This is the main benchmark in our section. ARC-Easy is only used later as a supporting general-capability check.",
        font_size=13,
        color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)


def slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title(slide, "SFT improved content generation; GRPO improved answer landing")
    slide.shapes.add_picture(str(ASSETS / "chart_gsm8k_core_metrics.png"), Inches(0.65), Inches(1.0), width=Inches(5.95))
    slide.shapes.add_picture(str(ASSETS / "chart_format_vs_content.png"), Inches(6.75), Inches(1.0), width=Inches(5.9))

    bullets = [
        "GRPO: best accuracy and strict_accuracy",
        "SFT: best content_accuracy",
        "Base: weakest on this diagnostic benchmark",
        "GRPO is better at structured answer delivery",
        "SFT is better at generating correct answer content",
        "These are gains on different stages of the same reasoning pipeline",
    ]
    add_bullets(slide, Inches(0.85), Inches(5.45), Inches(12.0), Inches(1.1), bullets, font_size=16)
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.35),
        Inches(12.0),
        Inches(0.28),
        "Key nuance: higher strict_accuracy does not automatically mean stronger content reasoning.",
        font_size=13,
        color=LIGHT_GRAY,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide)


def slide4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_bar(slide)
    add_title(slide, "Why GRPO wins on strict_accuracy, and what risks still remain")
    slide.shapes.add_picture(str(ASSETS / "chart_failure_modes.png"), Inches(0.6), Inches(1.0), width=Inches(5.9))
    slide.shapes.add_picture(str(ASSETS / "chart_format_coverage_diagnostic.png"), Inches(6.75), Inches(1.0), width=Inches(5.95))
    slide.shapes.add_picture(str(ASSETS / "chart_arc_easy_general_capability.png"), Inches(9.15), Inches(4.45), width=Inches(3.55))

    add_note_box(
        slide,
        Inches(0.75),
        Inches(5.2),
        Inches(4.0),
        Inches(1.2),
        "Why GRPO looks best on strict_accuracy",
        [
            "GRPO has explicit final answers in 296 / 300 samples.",
            "But inside the formatted subset, SFT has higher format_conditioned_accuracy.",
            "So GRPO gains mainly come from format coverage.",
        ],
    )
    add_note_box(
        slide,
        Inches(4.78),
        Inches(5.2),
        Inches(3.95),
        Inches(1.2),
        "Remaining risks",
        [
            "degeneration is still the biggest failure mode",
            "truncation is still very high at 1024 tokens",
            "format adherence and content correctness are still separated",
        ],
    )
    add_note_box(
        slide,
        Inches(9.08),
        Inches(5.92),
        Inches(3.45),
        Inches(0.95),
        "Supporting point",
        [
            "ARC-Easy shows no clear general-capability improvement.",
            "The gap is small, so this is a supporting observation, not the headline.",
        ],
    )
    add_footer(slide)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)

    prs.save(OUT)
    print(f"Saved PPT draft to: {OUT}")


if __name__ == "__main__":
    main()
